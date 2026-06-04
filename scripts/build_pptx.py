#!/usr/bin/env python3
"""Build inHunt presentation PPTX — v2."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Constants ──────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Colors
BG_DARK    = RGBColor(0x0F, 0x0F, 0x1A)
BG_CARD    = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT     = RGBColor(0x00, 0x7A, 0xFF)
ACCENT2    = RGBColor(0x34, 0xC7, 0x59)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0x8E, 0x8E, 0x93)
LIGHT_GRAY = RGBColor(0xC7, 0xC7, 0xCC)
DIM_WHITE  = RGBColor(0xAB, 0xAB, 0xB5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # Blank

TOTAL = 10  # 1 cover + 1 needs + 3 demo + 2 PE + 1 tech + 1 future + 1 thanks

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE,
             bold=False, align=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_multiline(slide, left, top, width, height, lines, font_size=14, color=WHITE,
                  font_name='Microsoft YaHei', line_spacing=1.3):
    """Add multi-line text with explicit line control."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_text in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line_text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(font_size * (line_spacing - 1) * 2)
    return txBox

def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_page_number(slide, num, total):
    add_text(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
             f'{num}/{total}', font_size=12, color=GRAY, align=PP_ALIGN.RIGHT)

def add_top_bar(slide):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Pt(4), ACCENT)

def add_section_header(slide, title, y=Inches(0.4)):
    add_text(slide, Inches(0.8), y, Inches(10), Inches(0.6),
             title, font_size=32, color=WHITE, bold=True)
    add_rect(slide, Inches(0.8), y + Inches(0.65), Inches(1.5), Pt(3), ACCENT)

# ═══════════════════════════════════════════════════════════
# SLIDE 1: Cover
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_top_bar(slide)

add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
         'inHunt 灵觅', font_size=54, color=WHITE, bold=True)
add_text(slide, Inches(1), Inches(2.7), Inches(11), Inches(0.8),
         '网文作者的 AI 灵感助手', font_size=28, color=ACCENT)
add_rect(slide, Inches(1), Inches(3.6), Inches(2), Pt(3), ACCENT)

add_text(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.5),
         '移动应用开发 · 课程大作业', font_size=18, color=DIM_WHITE)
add_text(slide, Inches(1), Inches(5.0), Inches(11), Inches(1.0),
         '答辩人：克里斯\n团队成员：[同学A · 开发]  [同学B · 产品创意]  [同学C]  [同学D]',
         font_size=16, color=GRAY)
add_text(slide, Inches(1), Inches(6.3), Inches(4), Inches(0.4),
         '2026 年 6 月', font_size=14, color=GRAY)

# ═══════════════════════════════════════════════════════════
# SLIDE 2: 需求分析
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_top_bar(slide)
add_section_header(slide, '需求分析')

pains = [
    ('💡 灵感碎片化', '灵感来了没地方记，随手写的纸条、备忘录散落各处'),
    ('📦 素材难管理', '角色、情节、世界观零散存储，写的时候找不到'),
    ('✍️ 写作缺辅助', '卡文时缺乏 AI 帮助，续写/润色需要反复切换工具'),
]
for i, (title, desc) in enumerate(pains):
    y = Inches(1.5 + i * 1.6)
    add_rounded_rect(slide, Inches(0.8), y, Inches(5.5), Inches(1.3), BG_CARD)
    add_text(slide, Inches(1.1), y + Inches(0.15), Inches(5), Inches(0.5),
             title, font_size=20, color=ACCENT, bold=True)
    add_text(slide, Inches(1.1), y + Inches(0.6), Inches(5), Inches(0.5),
             desc, font_size=15, color=DIM_WHITE)

add_text(slide, Inches(7), Inches(0.4), Inches(5.5), Inches(0.6),
         'inHunt 的解决方案', font_size=24, color=ACCENT2, bold=True)
add_rect(slide, Inches(7), Inches(1.0), Inches(1.5), Pt(3), ACCENT2)

solution_items = [
    '🎯 灵感捕捉 — AI 一键生成角色/情节/世界观灵感卡片',
    '📚 灵感记录 — 文件夹 + 标签分类，随时检索',
    '✍️ 故事管理 — AI 续写/润色 + 灵感便签 + 大纲 + 沉浸阅读',
    '⚙️ 灵活配置 — 用户自带 API Key，零后端部署',
]
for i, item in enumerate(solution_items):
    add_text(slide, Inches(7), Inches(1.5 + i * 0.7), Inches(5.8), Inches(0.6),
             item, font_size=16, color=WHITE)

add_rounded_rect(slide, Inches(7), Inches(4.5), Inches(5.8), Inches(1.0), ACCENT)
add_text(slide, Inches(7.3), Inches(4.55), Inches(5.2), Inches(0.9),
         '定位：从灵感捕捉到故事成稿的一站式 AI 写作助手',
         font_size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.4),
         '关键词：AI 灵感生成  ·  结构化素材管理  ·  AI 辅助写作  ·  Prompt Engineering',
         font_size=12, color=GRAY)
add_page_number(slide, 2, TOTAL)

# ═══════════════════════════════════════════════════════════
# SLIDE 3: Demo — 配置 & 灵感捕捉
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_top_bar(slide)
add_section_header(slide, '功能演示 ① — 配置 & 灵感捕捉')

# Left: description
add_text(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.4),
         '配置模型', font_size=20, color=ACCENT, bold=True)
add_multiline(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(1.5), [
    '▸ 填入 API Key（支持任意 OpenAI 兼容接口）',
    '▸ 选择模型（如 DeepSeek、GLM 等）',
    '▸ 调整温度、Top P 等参数',
    '▸ 设置写作偏好（文风、补充指令）',
], font_size=14, color=DIM_WHITE)

add_text(slide, Inches(0.8), Inches(3.6), Inches(5.5), Inches(0.4),
         '灵感捕捉', font_size=20, color=ACCENT, bold=True)
add_multiline(slide, Inches(0.8), Inches(4.1), Inches(5.5), Inches(2.0), [
    '▸ 四种预设模式：角色 / 情节 / 世界观 / 对话',
    '▸ 输入关键词，AI 生成 8 张灵感卡片',
    '▸ 每张卡片含标题 + 内容，可收藏到灵感记录',
    '▸ 支持自定义 prompt 补充描述',
], font_size=14, color=DIM_WHITE)

# Right: screenshot placeholder
add_rounded_rect(slide, Inches(7), Inches(1.3), Inches(5.8), Inches(5.5), BG_CARD)
add_text(slide, Inches(7.3), Inches(3.5), Inches(5.2), Inches(1.0),
         '[ 截图区域 ]\n配置页面 & 灵感捕捉页面',
         font_size=16, color=GRAY, align=PP_ALIGN.CENTER)

add_page_number(slide, 3, TOTAL)

# ═══════════════════════════════════════════════════════════
# SLIDE 4: Demo — 灵感记录 & 故事管理
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_top_bar(slide)
add_section_header(slide, '功能演示 ② — 灵感记录 & 故事管理')

# Left: screenshot placeholder
add_rounded_rect(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(5.5), BG_CARD)
add_text(slide, Inches(1.1), Inches(3.5), Inches(4.9), Inches(1.0),
         '[ 截图区域 ]\n灵感记录（文件夹/标签筛选）& 故事列表',
         font_size=16, color=GRAY, align=PP_ALIGN.CENTER)

# Right: description
add_text(slide, Inches(7), Inches(1.4), Inches(5.5), Inches(0.4),
         '灵感记录', font_size=20, color=ACCENT, bold=True)
add_multiline(slide, Inches(7), Inches(1.9), Inches(5.8), Inches(1.5), [
    '▸ 文件夹分类管理灵感素材',
    '▸ 标签筛选：角色 / 情节 / 世界观 / 对话',
    '▸ 编辑灵感内容，添加到故事',
], font_size=14, color=DIM_WHITE)

add_text(slide, Inches(7), Inches(3.5), Inches(5.5), Inches(0.4),
         '故事管理 — 编辑器', font_size=20, color=ACCENT, bold=True)
add_multiline(slide, Inches(7), Inches(4.0), Inches(5.8), Inches(2.8), [
    '▸ 创建故事，📌 关联灵感便签',
    '▸ 🌍 世界观设定面板（人物/关系/伏笔）',
    '▸ 章节大纲从正文自动解析',
    '▸ ✏️ AI 续写：自动带灵感+大纲+设定作为上下文',
    '▸ ✨ AI 润色：选中 → AI 改写 → 替换',
    '▸ 自动保存',
], font_size=14, color=DIM_WHITE)

add_page_number(slide, 4, TOTAL)

# ═══════════════════════════════════════════════════════════
# SLIDE 5: Demo — 沉浸阅读器
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_top_bar(slide)
add_section_header(slide, '功能演示 ③ — 沉浸阅读器')

# Large screenshot area
add_rounded_rect(slide, Inches(0.8), Inches(1.3), Inches(7.5), Inches(5.5), BG_CARD)
add_text(slide, Inches(1.1), Inches(3.5), Inches(6.9), Inches(1.0),
         '[ 截图区域 ]\n沉浸阅读器 — 羊皮纸风格 + 章节目录',
         font_size=16, color=GRAY, align=PP_ALIGN.CENTER)

# Right: features
add_text(slide, Inches(9), Inches(1.4), Inches(4), Inches(0.4),
         '阅读器特性', font_size=20, color=ACCENT, bold=True)
add_multiline(slide, Inches(9), Inches(2.0), Inches(3.8), Inches(4.0), [
    '▸ 羊皮纸暖底色 + 深棕文字',
    '▸ 隐藏 Tab 栏和标题栏',
    '▸ 章节目录弹窗跳转',
    '▸ 竖向滚动阅读',
    '▸ 滚动进度条',
    '▸ 点击切换 UI 显示',
], font_size=15, color=DIM_WHITE, line_spacing=1.5)

add_page_number(slide, 5, TOTAL)

# ═══════════════════════════════════════════════════════════
# SLIDE 6: PE ① — 灵感生成
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_top_bar(slide)
add_section_header(slide, 'Prompt Engineering ① — 结构化灵感生成')

add_text(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.4),
         'System Prompt（角色设定）', font_size=16, color=ACCENT, bold=True)
add_rounded_rect(slide, Inches(0.8), Inches(1.75), Inches(5.5), Inches(1.0), BG_CARD)
add_text(slide, Inches(1.0), Inches(1.8), Inches(5.1), Inches(0.9),
         '"你是一位网文角色设计师，擅长生成\n可直接发展为主角、配角或反派的\n人物种子。"',
         font_size=13, color=DIM_WHITE)

add_text(slide, Inches(0.8), Inches(3.0), Inches(5.5), Inches(0.4),
         'User Prompt（输出约束）', font_size=16, color=ACCENT, bold=True)
add_rounded_rect(slide, Inches(0.8), Inches(3.45), Inches(5.5), Inches(2.2), BG_CARD)
add_multiline(slide, Inches(1.0), Inches(3.5), Inches(5.1), Inches(2.1), [
    '• 返回 8 条灵感，不多不少',
    '• 每条必须有 title 和 content',
    '• content 控制在 65~85 个汉字',
    '• 不要解释、不要前言、不要 Markdown',
    '• 整个回复只能是严格 JSON 数组',
    '• 格式：[{"title":"...", "content":"..."}]',
], font_size=13, color=DIM_WHITE)

add_text(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(0.4),
         '工程容缺设计', font_size=16, color=ACCENT2, bold=True)
add_rounded_rect(slide, Inches(7), Inches(1.75), Inches(5.8), Inches(3.9), BG_CARD)
add_multiline(slide, Inches(7.2), Inches(1.85), Inches(5.4), Inches(3.7), [
    '问题：LLM 不一定遵守输出格式约束',
    '',
    '解法：多层容错解析器 (parseCards.ts)',
    '',
    '  ① 尝试 JSON.parse → 成功则直接用',
    '  ② 失败 → 按编号列表拆分 (1. xxx 2. xxx)',
    '  ③ 失败 → 按 Markdown 标题拆分',
    '  ④ 失败 → 按空行 / "标题：" 关键词拆分',
    '',
    '设计原则：',
    '  Prompt 约束是"理想路径"',
    '  工程容错是"底线保障"',
], font_size=14, color=WHITE)

add_rounded_rect(slide, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.8), ACCENT)
add_text(slide, Inches(1.0), Inches(6.25), Inches(11.3), Inches(0.7),
         '💡 要点：Prompt 约束 + 工程容缺 = 可靠的结构化输出',
         font_size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_page_number(slide, 6, TOTAL)

# ═══════════════════════════════════════════════════════════
# SLIDE 7: PE ② — 续写 vs 润色
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_top_bar(slide)
add_section_header(slide, 'Prompt Engineering ② — 续写 vs 润色的上下文工程')

# Left: 续写
add_text(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.4),
         '✏️ 续写模式', font_size=20, color=ACCENT, bold=True)
add_rounded_rect(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(3.5), BG_CARD)
add_multiline(slide, Inches(1.0), Inches(1.9), Inches(5.1), Inches(3.3), [
    'System Prompt:',
    '"你是小说写作助手。根据用户指令和',
    '上下文（正文、灵感、大纲、设定），',
    '生成合适的创作内容。"',
    '',
    '上下文拼装（自动注入）:',
    '  【当前正文（最近 3000 字）】',
    '  【故事设定】世界观 / 人物关系',
    '  【参考灵感】📌 便签内容',
    '  【大纲】章节标题列表',
    '  【文风要求】用户偏好',
    '',
    'Temperature: max(用户设定值, 0.8)',
    '  用户可调高，但保底 0.8 鼓励创造性',
], font_size=13, color=WHITE)

# Right: 润色
add_text(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(0.4),
         '✨ 润色模式', font_size=20, color=ACCENT2, bold=True)
add_rounded_rect(slide, Inches(7), Inches(1.8), Inches(5.8), Inches(3.5), BG_CARD)
add_multiline(slide, Inches(7.2), Inches(1.9), Inches(5.4), Inches(3.3), [
    'System Prompt:',
    '"【选中文本】是润色目标，',
    '【前文】【后文】只供参考上下文，',
    '不要修改它们。"',
    '',
    '上下文拼装（精确定位）:',
    '  【选中文本】← 润色目标',
    '  【前文】选区前 800 字（参考）',
    '  【后文】选区后 800 字（参考）',
    '  【故事设定】世界观 / 人物关系',
    '',
    'Temperature: 固定 0.7（保持稳定输出）',
    '  不受用户设置影响',
], font_size=13, color=WHITE)

# Bottom
add_rounded_rect(slide, Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.3), BG_CARD)
add_multiline(slide, Inches(1.0), Inches(5.7), Inches(11.3), Inches(1.1), [
    '用户偏好动态注入',
    '写作偏好（文风、补充指令）→ 写入 System Prompt → 用户修改后所有 AI 调用自动生效',
    '',
    '💡 续写和润色不是同一个 prompt 的变体，而是针对不同任务设计了不同的上下文结构、system prompt 和 temperature 策略',
], font_size=13, color=WHITE)
add_page_number(slide, 7, TOTAL)

# ═══════════════════════════════════════════════════════════
# SLIDE 8: 技术架构
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_top_bar(slide)
add_section_header(slide, '技术架构')

techs = [
    ('React Native', '跨平台移动框架', ACCENT),
    ('Expo 54', '开发/构建/部署工具链', ACCENT),
    ('TypeScript', '类型安全的开发语言', ACCENT),
    ('Expo Router', '基于文件的路由系统', ACCENT),
    ('AsyncStorage', '本地持久化存储', ACCENT2),
    ('LLM API', '用户自带 Key 的\nOpenAI 兼容接口', ACCENT2),
]
for i, (name, desc, color) in enumerate(techs):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + col * Inches(2.2)
    y = Inches(1.5) + row * Inches(1.8)
    add_rounded_rect(slide, x, y, Inches(2.0), Inches(1.5), BG_CARD)
    add_text(slide, x + Inches(0.1), y + Inches(0.15), Inches(1.8), Inches(0.5),
             name, font_size=16, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.1), y + Inches(0.65), Inches(1.8), Inches(0.7),
             desc, font_size=12, color=GRAY, align=PP_ALIGN.CENTER)

add_text(slide, Inches(7.5), Inches(1.3), Inches(5.5), Inches(0.4),
         '架构亮点', font_size=20, color=ACCENT, bold=True)

highlights = [
    ('纯前端架构', '所有数据 AsyncStorage 本地存储，\n无需后端服务器，零部署成本'),
    ('用户自带 API', '用户填入自己的 LLM API Key，\n支持任意 OpenAI 兼容接口'),
    ('模块化设计', '功能页面独立解耦，\nconfigStore 统一管理全局状态'),
]
for i, (title, desc) in enumerate(highlights):
    y = Inches(1.9 + i * 1.5)
    add_text(slide, Inches(7.5), y, Inches(5.5), Inches(0.4),
             f'▸ {title}', font_size=16, color=WHITE, bold=True)
    add_text(slide, Inches(7.8), y + Inches(0.4), Inches(5.2), Inches(0.8),
             desc, font_size=13, color=DIM_WHITE)

add_rounded_rect(slide, Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.5), BG_CARD)
add_text(slide, Inches(1.2), Inches(5.4), Inches(11), Inches(1.3),
         '项目规模    ·    20+ 源文件    ·    ~6000 行 TypeScript    ·    4 个功能模块    ·    完整 CRUD + AI 集成\n'
         '开源地址    ·    github.com/Chris-behind-door/Inspirational_Writer',
         font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_page_number(slide, 8, TOTAL)

# ═══════════════════════════════════════════════════════════
# SLIDE 9: 未来展望
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_top_bar(slide)
add_section_header(slide, '未来展望')

# Harness direction
add_text(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.4),
         'Harness Engineering 方向', font_size=20, color=ACCENT, bold=True)
add_multiline(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(1.8), [
    'Agent = Model + Harness — 通过构建控制层让 AI Agent 可靠可用',
    '',
    '▸ Sensors（反馈检测）：自动检测人物矛盾、时间线冲突、风格一致性',
    '▸ Feedback Loops（自检修正）：AI 写 → 自动检查 → 修正 → 用户确认',
    '▸ 从"单次 AI 调用"进化为"可控的创作 Agent"',
], font_size=15, color=WHITE)

# Other future items
add_text(slide, Inches(0.8), Inches(3.8), Inches(11), Inches(0.4),
         '功能规划', font_size=20, color=ACCENT, bold=True)

future_items = [
    ('🔍 全文搜索', '故事正文关键词搜索\n结果高亮 + 位置跳转'),
    ('☁️ 多端同步', '云端备份\n跨设备无缝切换'),
    ('🌐 多语言支持', '中英文界面切换\n服务海外创作者'),
]
for i, (title, desc) in enumerate(future_items):
    x = Inches(0.8) + i * Inches(4.0)
    add_rounded_rect(slide, x, Inches(4.4), Inches(3.7), Inches(1.8), BG_CARD)
    add_text(slide, x + Inches(0.15), Inches(4.55), Inches(3.4), Inches(0.4),
             title, font_size=16, color=ACCENT, bold=True)
    add_text(slide, x + Inches(0.15), Inches(5.0), Inches(3.4), Inches(1.0),
             desc, font_size=13, color=DIM_WHITE)

add_page_number(slide, 9, TOTAL)

# ═══════════════════════════════════════════════════════════
# SLIDE 10: 致谢 & Q&A
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_top_bar(slide)

add_text(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(1.0),
         '致谢', font_size=40, color=WHITE, bold=True)
add_rect(slide, Inches(1), Inches(2.5), Inches(2), Pt(3), ACCENT)

add_multiline(slide, Inches(1), Inches(3.0), Inches(11), Inches(2.0), [
    '感谢课程老师的悉心指导',
    '感谢团队成员的协作与支持',
    '',
    '[ 可补充其他致谢对象 ]',
], font_size=18, color=DIM_WHITE, line_spacing=1.5)

# Q&A block
add_rounded_rect(slide, Inches(3.5), Inches(5.2), Inches(6.3), Inches(1.2), ACCENT)
add_text(slide, Inches(3.7), Inches(5.25), Inches(5.9), Inches(1.1),
         '谢谢！欢迎提问 🎤',
         font_size=32, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_page_number(slide, 10, TOTAL)

# ── Save ───────────────────────────────────────────────────
output = os.path.expanduser('~/下载/inHunt_演示PPT.pptx')
prs.save(output)
print(f'Saved to {output}  ({os.path.getsize(output)} bytes)')
